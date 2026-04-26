#!/usr/bin/env python3
"""Generate the B-BOT final presentation deck."""

from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
PRESENTATION_DIR = ROOT / "Report" / "presentation"
DEPS = PRESENTATION_DIR / ".deps"
if DEPS.exists():
    sys.path.insert(0, str(DEPS))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image


OUT = PRESENTATION_DIR / "B-BOT_final_presentation.pptx"

FIG = ROOT / "Report" / "figures"
APP = ROOT / "Report" / "appendices"

COLORS = {
    "bg": RGBColor(248, 250, 252),
    "ink": RGBColor(22, 32, 43),
    "muted": RGBColor(85, 97, 110),
    "blue": RGBColor(32, 91, 167),
    "teal": RGBColor(0, 130, 130),
    "orange": RGBColor(209, 111, 31),
    "green": RGBColor(63, 133, 82),
    "red": RGBColor(185, 64, 56),
    "line": RGBColor(210, 218, 228),
    "white": RGBColor(255, 255, 255),
}


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["bg"]
    return slide


def add_text(slide, text, x, y, w, h, size=24, bold=False, color="ink", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.55, 0.35, 12.0, 0.55, size=30, bold=True, color="ink")
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.95, 11.5, 0.32, size=13, color="muted")
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.28), Inches(2.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["teal"]
    line.line.color.rgb = COLORS["teal"]


def add_bullets(slide, items, x, y, w, h, size=18, color="ink", level_gap=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.space_after = Pt(7)
        p.font.name = "Aptos"
        p.font.size = Pt(max(size - level * level_gap, 11))
        p.font.color.rgb = COLORS[color]
    return box


def add_card(slide, x, y, w, h, title, body, accent="blue"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    add_text(slide, title, x + 0.18, y + 0.15, w - 0.36, 0.32, size=15, bold=True, color=accent)
    if isinstance(body, list):
        add_bullets(slide, body, x + 0.18, y + 0.57, w - 0.35, h - 0.65, size=13, color="ink")
    else:
        add_text(slide, body, x + 0.18, y + 0.58, w - 0.35, h - 0.65, size=13, color="ink")


def add_image_fit(slide, path, x, y, w, h, border=True):
    path = Path(path)
    with Image.open(path) as img:
        iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        width = w
        height = w / img_ratio
        left = x
        top = y + (h - height) / 2
    else:
        height = h
        width = h * img_ratio
        left = x + (w - width) / 2
        top = y
    pic = slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    if border:
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        rect.fill.background()
        rect.line.color.rgb = COLORS["line"]
        rect.line.width = Pt(1)
    return pic


def add_footer(slide, idx):
    add_text(slide, f"B-BOT Final Presentation | {idx}", 10.9, 7.05, 1.75, 0.22, size=8, color="muted", align=PP_ALIGN.RIGHT)


def connector(slide, x1, y1, x2, y2, color="muted"):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = COLORS[color]
    line.line.width = Pt(1.5)


def generate():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1
    slide = blank_slide(prs)
    add_text(slide, "B-BOT", 0.7, 0.85, 6.0, 0.8, size=48, bold=True, color="ink")
    add_text(slide, "A WiFi-Enabled Self-Balancing Wheel-Legged Robot\nwith ROS 2 Vision Teleoperation", 0.72, 1.78, 8.8, 0.85, size=24, color="blue")
    add_text(slide, "Botao Su | Third Year Individual Project | University of Manchester", 0.75, 6.35, 10.5, 0.3, size=16, color="muted")
    add_card(slide, 8.5, 1.15, 3.7, 3.9, "Core Thesis", [
        "Balance-critical feedback stays on the ESP32.",
        "WiFi, Xbox and ROS 2 vision send supervisory targets only.",
        "Safety arbitration is measured, not assumed.",
    ], accent="teal")
    add_footer(slide, 1)

    # 2
    slide = blank_slide(prs)
    add_title(slide, "Project Problem", "An unstable robot should not depend on non-deterministic communication for balance.")
    add_card(slide, 0.8, 1.75, 3.6, 3.7, "Wheel-Legged Challenge", [
        "Changing leg length changes the balance dynamics.",
        "Pitch, wheel torque and leg geometry are coupled.",
        "A fixed simple controller is fragile across configurations.",
    ], accent="blue")
    add_card(slide, 4.85, 1.75, 3.6, 3.7, "Teleoperation Challenge", [
        "WiFi and vision have variable latency.",
        "Human input sources can conflict.",
        "Unsafe stale commands must be cleared predictably.",
    ], accent="orange")
    add_card(slide, 8.9, 1.75, 3.6, 3.7, "Design Rule", [
        "ESP32 owns the stabilising loop.",
        "Remote inputs are target requests.",
        "Watchdogs and arbitration enforce safe behaviour.",
    ], accent="teal")
    add_footer(slide, 2)

    # 3
    slide = blank_slide(prs)
    add_title(slide, "Objectives", "The report and evaluation are organised around three testable objectives.")
    add_card(slide, 0.8, 1.6, 3.7, 4.3, "O1 Embedded Balance", [
        "LQR/PID/VMC balance control.",
        "Static stability and disturbance recovery.",
        "4 ms local control-loop timing.",
    ], accent="blue")
    add_card(slide, 4.8, 1.6, 3.7, 4.3, "O2 Vision + WiFi", [
        "ROS 2 camera and MediaPipe bridge.",
        "WiFi TCP command-entry path.",
        "Measured supervisory latency.",
    ], accent="teal")
    add_card(slide, 8.8, 1.6, 3.7, 4.3, "O3 Safe Arbitration", [
        "BLE, UART, TCP and vision command sources.",
        "Watchdog stop behaviour.",
        "Conflict suppression and dry-run/stunt gates.",
    ], accent="orange")
    add_footer(slide, 3)

    # 4 architecture
    slide = blank_slide(prs)
    add_title(slide, "System Architecture", "Time-critical balance is separated from supervisory command generation.")
    layers = [
        ("Host PC / ROS 2 / MediaPipe", 0.9, 1.55, 3.4, 0.8, "teal"),
        ("Xbox BLE / UART / WiFi TCP", 5.0, 1.55, 3.4, 0.8, "orange"),
        ("ESP32 Firmware\nFreeRTOS tasks, parser, watchdogs", 3.2, 3.0, 4.6, 0.95, "blue"),
        ("Local Balance Control\nIMU + motor feedback, LQR/PID/VMC", 3.2, 4.55, 4.6, 0.95, "green"),
        ("Physical Robot\nwheels, legs, CAN motors, battery", 3.2, 6.0, 4.6, 0.75, "muted"),
    ]
    for text, x, y, w, h, color in layers:
        add_card(slide, x, y, w, h, text, "", accent=color)
    connector(slide, 2.6, 2.35, 4.2, 3.0, "muted")
    connector(slide, 6.7, 2.35, 6.0, 3.0, "muted")
    connector(slide, 5.5, 3.95, 5.5, 4.55, "muted")
    connector(slide, 5.5, 5.5, 5.5, 6.0, "muted")
    add_bullets(slide, [
        "Remote layers update targets, not stabilising feedback.",
        "Command loss or conflict is handled before it can become persistent motion.",
        "This boundary is tested by E5, E6, E8, E10, E11 and O3-S.",
    ], 8.6, 3.05, 3.8, 2.5, size=16)
    add_footer(slide, 4)

    # 5 hardware
    slide = blank_slide(prs)
    add_title(slide, "Hardware and Firmware Implementation", "A lightweight embedded stack built around ESP32, IMU feedback and CAN motor control.")
    add_image_fit(slide, APP / "Appendix_H_Hardware_and_Control_Evidence" / "PCB_Layout.png", 0.8, 1.55, 5.2, 4.8)
    add_card(slide, 6.35, 1.6, 5.8, 4.55, "Implementation Evidence", [
        "Custom ESP32 controller PCB and wheel-legged platform.",
        "MPU6050 attitude feedback and CAN-connected motors.",
        "FreeRTOS tasks for CAN, IMU, motor output, balance, BLE, UART and WiFi.",
        "Public software evidence snapshot: Realsbt/Y3-PROJECT.",
    ], accent="blue")
    add_footer(slide, 5)

    # 6 control
    slide = blank_slide(prs)
    add_title(slide, "Balance Control Hierarchy", "LQR handles coupled sagittal balance; PID and VMC handle local regulation and leg torque mapping.")
    boxes = [
        ("State Estimate\npitch, pitch rate,\nwheel speed, leg length", 0.8, 2.0, "blue"),
        ("Gain-Scheduled LQR\nsagittal balance", 3.6, 2.0, "teal"),
        ("PID Loops\nyaw, roll, leg length", 6.4, 2.0, "orange"),
        ("VMC Mapping\nvirtual leg forces\nto joint torques", 9.2, 2.0, "green"),
        ("Motor Commands\nwheels + joints", 5.0, 4.75, "blue"),
    ]
    for text, x, y, color in boxes:
        add_card(slide, x, y, 2.35, 1.3, text, "", accent=color)
    connector(slide, 3.15, 2.65, 3.6, 2.65)
    connector(slide, 5.95, 2.65, 6.4, 2.65)
    connector(slide, 8.75, 2.65, 9.2, 2.65)
    connector(slide, 10.35, 3.3, 6.15, 4.75)
    add_bullets(slide, [
        "Chosen for ESP32 feasibility and interpretable control structure.",
        "E3 and E9 test whether leg-length-dependent behaviour matters.",
        "Ablation is used to justify complexity rather than just describe it.",
    ], 0.95, 5.35, 4.3, 1.25, size=15)
    add_footer(slide, 6)

    # 7 command safety
    slide = blank_slide(prs)
    add_title(slide, "Command and Safety Architecture", "Multiple command sources are allowed, but persistent unsafe motion is not.")
    add_card(slide, 0.75, 1.55, 2.45, 1.45, "Command Sources", ["Xbox BLE", "UART queue", "WiFi TCP", "ROS 2 vision"], accent="orange")
    add_card(slide, 3.75, 1.55, 2.7, 1.45, "Parser + Queue", ["Shared command grammar", "Queue start/pause/stop", "Direct DRIVE/YAWRATE"], accent="blue")
    add_card(slide, 7.0, 1.55, 2.55, 1.45, "Safety Gates", ["Watchdogs", "BLE disable", "dry_run", "stunt_armed"], accent="teal")
    add_card(slide, 10.1, 1.55, 2.2, 1.45, "Robot Targets", ["speed", "yaw rate", "leg action", "stop"], accent="green")
    connector(slide, 3.2, 2.3, 3.75, 2.3)
    connector(slide, 6.45, 2.3, 7.0, 2.3)
    connector(slide, 9.55, 2.3, 10.1, 2.3)
    add_bullets(slide, [
        "Direct commands expire if not refreshed.",
        "Queue execution suppresses conflicting direct drive input.",
        "Vision can be previewed in dry-run before real transmission.",
        "O3-S verifies BLE/PC and queue/direct-command arbitration cases.",
    ], 1.0, 4.05, 10.9, 1.8, size=18)
    add_footer(slide, 7)

    # 8 eval strategy
    slide = blank_slide(prs)
    add_title(slide, "Evaluation Strategy", "Eleven measured experiments plus O3 arbitration supplement.")
    add_card(slide, 0.75, 1.6, 3.7, 4.4, "O1 Evidence", [
        "E1 static balance: 0.291 deg pitch RMS.",
        "E2 recovery: 9/10 forward, 10/10 backward.",
        "E8 loop timing: mean 3.9998 ms, median 4.000 ms.",
        "E9 controller ablation: FULL vs reduced variants.",
    ], accent="blue")
    add_card(slide, 4.8, 1.6, 3.7, 4.4, "O2 Evidence", [
        "E5 TCP ACK median 37.41 ms.",
        "E7 camera bring-up.",
        "E10 live gesture matrix and retained failures.",
        "E11 bridge ACK median 66.13 ms.",
    ], accent="teal")
    add_card(slide, 8.85, 1.6, 3.7, 4.4, "O3 Evidence", [
        "E6 watchdog and disconnect fault injection.",
        "dry_run and stunt gate behaviour.",
        "O3-S: 5/5 BLE gate and 5/5 queue/direct suppression.",
        "Not claimed as exhaustive over every race condition.",
    ], accent="orange")
    add_footer(slide, 8)

    # 9 O1
    slide = blank_slide(prs)
    add_title(slide, "O1 Results: Static Balance and Recovery", "Physical results support usable balance while exposing limits.")
    add_image_fit(slide, FIG / "measured" / "e1_static_balance_drift.png", 0.65, 1.55, 5.6, 3.8)
    add_image_fit(slide, FIG / "measured" / "e2_disturbance_recovery.png", 6.65, 1.55, 5.8, 3.8)
    add_bullets(slide, [
        "Static balance: 0.291 deg pitch RMS over three 60 s trials.",
        "Impulse recovery: 9/10 forward and 10/10 backward trials.",
        "The failed/protected forward case is retained as evidence of a real tuning/saturation limit.",
    ], 0.85, 5.7, 11.5, 0.9, size=15)
    add_footer(slide, 9)

    # 10 E8
    slide = blank_slide(prs)
    add_title(slide, "O1 Timing Boundary", "The local loop is centred on 4 ms, but timing outliers are real.")
    add_image_fit(slide, FIG / "e8_loop_jitter_15000_2026-04-24.png", 0.75, 1.45, 6.4, 4.75)
    add_card(slide, 7.55, 1.75, 4.4, 3.9, "Measured E8 Result", [
        "15,000 samples.",
        "Mean period: 3.9998 ms.",
        "Median period: 4.000 ms.",
        "p95: 4.300 ms.",
        "Maximum observed period: 53.365 ms.",
        "Conclusion: WiFi/vision must remain outside the stabilising loop.",
    ], accent="teal")
    add_footer(slide, 10)

    # 11 E9
    slide = blank_slide(prs)
    add_title(slide, "O1 Baseline: Controller Ablation", "E9 tests whether added control complexity is justified.")
    add_image_fit(slide, FIG / "measured" / "e9_controller_ablation.png", 0.75, 1.35, 6.4, 4.95)
    add_card(slide, 7.5, 1.6, 4.75, 4.25, "FULL vs FIXED_LQR", [
        "FULL: 10/10 recovered.",
        "FIXED_LQR: 9/10 recovered, one protected/failed trial.",
        "Response metric: 1.049 s to 0.826 s, 21.3% reduction.",
        "Peak pitch: 9.85 deg to 8.09 deg, 17.9% reduction.",
        "Interpretation: engineering ablation evidence, not proof of global optimality.",
    ], accent="blue")
    add_footer(slide, 11)

    # 12 O2
    slide = blank_slide(prs)
    add_title(slide, "O2 Results: WiFi and Vision Teleoperation", "Communication and vision are responsive enough for supervisory commands.")
    add_image_fit(slide, FIG / "e10_vision_confusion_live_2026-04-24.png", 0.65, 1.45, 5.7, 4.4)
    add_image_fit(slide, FIG / "e11_vision_bridge_ack_latency_2026-04-24.png", 6.85, 1.45, 5.65, 4.4)
    add_bullets(slide, [
        "TCP command-entry ACK: median 37.41 ms, p95 88.31 ms across 300 samples.",
        "Vision live matrix: clean 6/6 command classes; 85.3% selected-frame label accuracy.",
        "Vision bridge-to-ESP32 ACK: median 66.13 ms across 71 safe commands.",
    ], 0.85, 6.0, 11.6, 0.7, size=13)
    add_footer(slide, 12)

    # 13 O3
    slide = blank_slide(prs)
    add_title(slide, "O3 Results: Safety and Arbitration", "Safety behaviour is measured on the main command paths.")
    add_image_fit(slide, FIG / "measured" / "e6_watchdog_fault_injection.png", 0.7, 1.45, 6.2, 4.8)
    add_card(slide, 7.35, 1.65, 4.8, 4.2, "Measured Safety Evidence", [
        "Direct command watchdog cleared commands in 10/10 trials.",
        "TCP socket close detected in 10/10 trials.",
        "TCP idle produced full-stop in 10/10 trials.",
        "O3-S: BLE/PC gate 5/5 and queue/direct suppression 5/5.",
        "Scope is tested main paths, not every possible race condition.",
    ], accent="orange")
    add_footer(slide, 13)

    # 14 demo
    slide = blank_slide(prs)
    add_title(slide, "Live Demonstration Plan", "Keep the demo conservative: show architecture, not maximum aggression.")
    add_card(slide, 0.8, 1.55, 3.5, 4.8, "Demo A: Local Balance", [
        "Robot already powered and calibrated.",
        "Show stable balance briefly.",
        "Small Xbox or keyboard command.",
        "Return to zero motion.",
    ], accent="blue")
    add_card(slide, 4.9, 1.55, 3.5, 4.8, "Demo B: Supervisory Command", [
        "Show terminal command path.",
        "Explain DRIVE/YAWRATE as targets.",
        "Show watchdog or stop behaviour conservatively.",
    ], accent="teal")
    add_card(slide, 9.0, 1.55, 3.5, 4.8, "Demo C: Vision Dry-Run", [
        "Run vision bridge in dry_run=true.",
        "Show gesture classified into command.",
        "No real robot motion unless already rehearsed.",
    ], accent="orange")
    add_footer(slide, 14)

    # 15 limitations
    slide = blank_slide(prs)
    add_title(slide, "Limitations and Future Work", "The report is strongest when the boundaries are explicit.")
    add_card(slide, 0.8, 1.6, 5.4, 4.6, "Current Limitations", [
        "Manual gain tuning and manually applied disturbances.",
        "Small-sample physical trials.",
        "IMU vibration, mechanical compliance and actuator saturation.",
        "WiFi/vision latency prevents use as balance feedback.",
    ], accent="red")
    add_card(slide, 6.95, 1.6, 5.4, 4.6, "Next Technical Steps", [
        "Better logging, event markers and parameter versioning.",
        "Repeatable gain sweeps and larger E2/E9 datasets.",
        "Improved state estimation and vibration isolation.",
        "Onboard companion computer for ROS 2 vision, while ESP32 still owns balance.",
    ], accent="green")
    add_footer(slide, 15)

    # 16 conclusion
    slide = blank_slide(prs)
    add_title(slide, "Conclusion", "B-BOT demonstrates a practical embedded balance architecture with safe supervisory teleoperation.")
    add_card(slide, 0.9, 1.55, 11.3, 4.6, "Main Contribution", [
        "ESP32 owns the balance-critical feedback loop.",
        "LQR/PID/VMC control is implemented and physically evaluated.",
        "WiFi, Xbox and ROS 2 vision are integrated as supervisory command sources.",
        "Watchdogs, dry-run, stunt gates and command arbitration are measured.",
        "The final system is not perfect, but its evidence chain is objective and traceable.",
    ], accent="teal")
    add_text(slide, "Thank you. I am happy to answer questions.", 2.1, 6.35, 9.2, 0.35, size=24, bold=True, color="blue", align=PP_ALIGN.CENTER)
    add_footer(slide, 16)

    # 17 backup
    slide = blank_slide(prs)
    add_title(slide, "Backup: Key Numbers for Q&A", "Use this slide only if asked for evidence details.")
    add_card(slide, 0.8, 1.55, 3.7, 4.7, "Balance", [
        "E1 pitch RMS: 0.291 deg.",
        "E2 recovery: 9/10 forward, 10/10 backward.",
        "E8 mean/median: 3.9998 ms / 4.000 ms.",
        "E8 max outlier: 53.365 ms.",
    ], accent="blue")
    add_card(slide, 4.8, 1.55, 3.7, 4.7, "Control Baseline", [
        "E9 FULL: 10/10.",
        "E9 FIXED_LQR: 9/10.",
        "Response metric: 21.3% lower.",
        "Peak pitch: 17.9% lower.",
    ], accent="teal")
    add_card(slide, 8.8, 1.55, 3.7, 4.7, "Teleoperation + Safety", [
        "E5 TCP median/p95: 37.41 / 88.31 ms.",
        "E11 ACK median: 66.13 ms.",
        "E10 selected-frame accuracy: 85.3%.",
        "E6 and O3-S main safety tests passed.",
    ], accent="orange")
    add_footer(slide, 17)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    generate()
